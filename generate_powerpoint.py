"""
Time Table Management System - PowerPoint Presentation Generator
Creates a professional presentation with all system details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import sqlite3
from datetime import datetime

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)  # Dark blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(66, 133, 244)
    p.alignment = PP_ALIGN.CENTER
    
    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = f"Generated: {datetime.now().strftime('%d-%m-%Y')}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_type="bullet"):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(66, 133, 244)
    header_shape.line.color.rgb = RGBColor(66, 133, 244)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    content_box.word_wrap = True
    
    return slide, content_box

def add_table_slide(prs, title, data, col_widths=None):
    """Add a slide with a table"""
    slide, content_box = add_content_slide(prs, title)
    
    # Create table
    rows, cols = len(data), len(data[0])
    left = Inches(0.7)
    top = Inches(1.3)
    width = Inches(8.6)
    height = Inches(5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table_shape.columns[i].width = Inches(w)
    
    # Fill data
    for i, row in enumerate(data):
        for j, cell_data in enumerate(row):
            cell = table_shape.cell(i, j)
            cell.text = str(cell_data)
            
            # Format header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(52, 73, 94)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.bold = True
                        run.font.size = Pt(11)
            else:
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(236, 240, 241)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
    
    return slide

def generate_presentation():
    """Generate the PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    add_title_slide(prs, "TIME TABLE MANAGEMENT", "System Documentation & Technical Overview")
    
    # SLIDE 2: Overview
    slide, content_box = add_content_slide(prs, "System Overview")
    tf = content_box.text_frame
    tf.clear()
    
    items = [
        ("Comprehensive Solution", "Manages educational institution schedules efficiently"),
        ("Database-Driven", "Built on SQLite with relational database design"),
        ("GUI Interface", "Professional Tkinter-based user interface"),
        ("Easy to Use", "Intuitive menu-driven system for staff and administrators"),
        ("Scalable", "Designed for growing data and multiple departments"),
        ("Real-time Management", "Live schedule updates and conflict detection"),
    ]
    
    for i, (title, desc) in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{title}: {desc}"
        p.level = 0
        p.font.size = Pt(14)
        p.space_before = Pt(8)
        if i == 0:
            p.text = f"{title}: {desc}"
    
    # SLIDE 3: Database Architecture
    slide, content_box = add_content_slide(prs, "Database Architecture")
    tf = content_box.text_frame
    tf.clear()
    
    items = [
        "6 Core Tables with Relational Design",
        "Automatic Primary & Foreign Keys",
        "Data Integrity & Constraints",
        "Timestamp Tracking (created_at, updated_at)",
        "Unique Constraints for Data Validation",
        "Indexed Fields for Fast Queries",
        "Support for Complex JOINs",
        "Scalable for 1000s of entries",
    ]
    
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"✓ {item}"
        p.level = 0
        p.font.size = Pt(14)
        p.space_before = Pt(5)
    
    # SLIDE 4: Table Structure Overview
    table_data = [
        ["Table Name", "Purpose", "Key Fields"],
        ["teachers", "Faculty information", "ID, Name, Department, Email"],
        ["classrooms", "Room details", "ID, Name, Capacity, Building"],
        ["subjects", "Course information", "ID, Code, Name, Credits"],
        ["time_slots", "Available periods", "ID, Start, End, Duration"],
        ["days", "Days of week", "ID, Name, Order"],
        ["timetable", "Main mapping table", "All Foreign Keys, Status"],
    ]
    add_table_slide(prs, "Table Structures", table_data, [2.2, 2.2, 2.2])
    
    # SLIDE 5: Teachers Table Details
    table_data = [
        ["Field Name", "Type", "Description"],
        ["teacher_id", "INTEGER", "Primary Key (Auto-Increment)"],
        ["teacher_name", "TEXT", "Name of Teacher (NOT NULL)"],
        ["department", "TEXT", "Department (NOT NULL)"],
        ["specialization", "TEXT", "Area of expertise"],
        ["contact_number", "TEXT", "Phone number"],
        ["email", "TEXT", "Email (UNIQUE)"],
    ]
    add_table_slide(prs, "Teachers Table Schema", table_data, [1.8, 1.5, 2.9])
    
    # SLIDE 6: Classrooms & Subjects Tables
    table_data = [
        ["Field Name", "Type", "Description"],
        ["classroom_id", "INTEGER", "Primary Key"],
        ["class_name", "TEXT", "Name (UNIQUE)"],
        ["capacity", "INTEGER", "Max students"],
        ["room_number", "TEXT", "Room ID"],
        ["building", "TEXT", "Building name"],
    ]
    add_table_slide(prs, "Classrooms Table Schema", table_data, [1.8, 1.5, 2.9])
    
    # SLIDE 7: TimeTable Schema (Main)
    table_data = [
        ["Field Name", "Type"],
        ["timetable_id", "INTEGER (PK, Auto-Increment)"],
        ["classroom_id", "INTEGER (FK)"],
        ["teacher_id", "INTEGER (FK)"],
        ["subject_id", "INTEGER (FK)"],
        ["day_id", "INTEGER (FK)"],
        ["slot_id", "INTEGER (FK)"],
        ["academic_year", "TEXT"],
        ["semester", "INTEGER"],
        ["session_type", "TEXT (Lecture/Lab/Practical)"],
        ["status", "TEXT (Active/Inactive)"],
    ]
    add_table_slide(prs, "TimeTable Schema (Central Table)", table_data, [2.5, 3.6])
    
    # SLIDE 8: System Statistics
    try:
        conn = sqlite3.connect('timetable.db')
        cursor = conn.cursor()
        
        cursor.execute("SELECT COUNT(*) FROM teachers")
        teachers_count = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM classrooms")
        classrooms_count = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM subjects")
        subjects_count = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM time_slots")
        slots_count = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM timetable")
        timetable_count = cursor.fetchone()[0]
        
        conn.close()
        
        slide, content_box = add_content_slide(prs, "System Statistics")
        
        # Add statistics as text
        tf = content_box.text_frame
        tf.clear()
        
        stats = [
            f"Total Teachers: {teachers_count}",
            f"Total Classrooms: {classrooms_count}",
            f"Total Subjects: {subjects_count}",
            f"Total Time Slots: {slots_count}",
            f"Total Schedule Entries: {timetable_count}",
        ]
        
        for i, stat in enumerate(stats):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"📊 {stat}"
            p.level = 0
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_before = Pt(15)
            p.font.color.rgb = RGBColor(66, 133, 244)
        
    except Exception as e:
        print(f"Error reading statistics: {e}")
    
    # SLIDE 9: SQL Query Examples - SELECT
    slide, content_box = add_content_slide(prs, "SQL Queries - SELECT Examples")
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Get All Teachers:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)
    
    p = tf.add_paragraph()
    p.text = "SELECT * FROM teachers ORDER BY teacher_id;"
    p.level = 0
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(231, 76, 60)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Get Time Table with Joins:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)
    
    p = tf.add_paragraph()
    p.text = "SELECT * FROM timetable t JOIN teachers JOIN classrooms..."
    p.level = 0
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(231, 76, 60)
    
    # SLIDE 10: SQL Query Examples - INSERT/UPDATE/DELETE
    slide, content_box = add_content_slide(prs, "SQL Queries - CRUD Operations")
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "INSERT New Entry:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)
    
    p = tf.add_paragraph()
    p.text = "INSERT INTO timetable (...) VALUES (...);"
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "UPDATE Entry:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "UPDATE timetable SET teacher_id = ? WHERE id = ?;"
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "DELETE Entry:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "DELETE FROM timetable WHERE timetable_id = ?;"
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    
    # SLIDE 11: GUI Features
    slide, content_box = add_content_slide(prs, "GUI Features & Interface")
    tf = content_box.text_frame
    tf.clear()
    
    features = [
        "Professional Tkinter Interface",
        "Login Window with Authentication",
        "Dashboard with Main Menu",
        "Data Tables with Tree View",
        "Add/Edit/Delete Operations",
        "Search and Filter Capabilities",
        "Color-Coded Status Display",
        "Real-time Data Updates",
        "Export to Text Files",
        "User-Friendly Navigation",
    ]
    
    for i, feature in enumerate(features):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"✓ {feature}"
        p.level = 0
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # SLIDE 12: Advanced Queries
    slide, content_box = add_content_slide(prs, "Advanced Query Examples")
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Count Classes by Day:"
    p.font.size = Pt(12)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "SELECT day_name, COUNT(*) FROM timetable GROUP BY day_id"
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "Get Teacher Schedule:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "SELECT * FROM timetable WHERE teacher_id = ?"
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "Available Slots:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "SELECT * FROM time_slots NOT IN (SELECT slot_id...)"
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    
    # SLIDE 13: Sample Data Display
    sample_data = [
        ["Teacher ID", "Name", "Department"],
        ["1", "Dr. John Smith", "Computer Science"],
        ["2", "Prof. Sarah Jones", "Mathematics"],
        ["3", "Dr. Mike Williams", "Physics"],
        ["4", "Ms. Emily Brown", "Chemistry"],
        ["5", "Dr. David Lee", "Biology"],
    ]
    add_table_slide(prs, "Sample Teachers Data", sample_data, [1.8, 2.5, 2.3])
    
    # SLIDE 14: Sample Classrooms
    sample_data = [
        ["ID", "Class Name", "Capacity", "Room #"],
        ["1", "Class A1", "40", "101"],
        ["2", "Class A2", "35", "102"],
        ["3", "Class B1", "45", "201"],
        ["4", "Class B2", "40", "202"],
        ["5", "Lab 1", "30", "301"],
    ]
    add_table_slide(prs, "Sample Classrooms Data", sample_data, [1, 2, 1.5, 1.5])
    
    # SLIDE 15: Sample Subjects
    sample_data = [
        ["ID", "Code", "Subject Name", "Credits"],
        ["1", "CS101", "Programming Basics", "3"],
        ["2", "CS102", "Data Structures", "4"],
        ["3", "MTH201", "Calculus I", "3"],
        ["4", "PHY101", "Physics I", "3"],
        ["5", "CHM101", "Chemistry I", "3"],
    ]
    add_table_slide(prs, "Sample Subjects Data", sample_data, [0.8, 1.3, 2.3, 1])
    
    # SLIDE 16: Sample Time Slots
    sample_data = [
        ["ID", "Slot Name", "Start Time", "End Time"],
        ["1", "Slot 1", "08:00", "09:00"],
        ["2", "Slot 2", "09:15", "10:15"],
        ["3", "Slot 3", "10:30", "11:30"],
        ["4", "Slot 4", "11:45", "12:45"],
        ["5", "Slot 5", "13:30", "14:30"],
    ]
    add_table_slide(prs, "Sample Time Slots", sample_data, [1, 1.5, 1.5, 1.5])
    
    # SLIDE 17: Key Features & Benefits
    slide, content_box = add_content_slide(prs, "Key Features & Benefits")
    tf = content_box.text_frame
    tf.clear()
    
    features = [
        "Automated Schedule Management",
        "No Time Slot Conflicts",
        "Quick Classroom/Teacher Search",
        "Department-wise Organization",
        "Semester/Year-based Planning",
        "Session Type Support (Lecture/Lab/Practical)",
        "Status Tracking (Active/Inactive)",
        "Historical Data Retention",
        "Export & Reporting Capability",
        "Easy Data Backup & Recovery",
    ]
    
    for i, feature in enumerate(features):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {feature}"
        p.level = 0
        p.font.size = Pt(13)
        p.space_before = Pt(5)
    
    # SLIDE 18: Technology Stack
    slide, content_box = add_content_slide(prs, "Technology Stack")
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Backend:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 133, 244)
    
    for tech in ["Python 3.14.0", "SQLite Database", "SQL Queries & Joins"]:
        p = tf.add_paragraph()
        p.text = f"  ✓ {tech}"
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Frontend:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 133, 244)
    p.space_before = Pt(10)
    
    for tech in ["Tkinter GUI Framework", "Tree View Widget", "Dialog Boxes & Pop-ups"]:
        p = tf.add_paragraph()
        p.text = f"  ✓ {tech}"
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Tools & Utilities:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 133, 244)
    p.space_before = Pt(10)
    
    for tech in ["Batch Launcher Scripts", "PDF Report Generation", "Data Export Functions"]:
        p = tf.add_paragraph()
        p.text = f"  ✓ {tech}"
        p.font.size = Pt(12)
    
    # SLIDE 19: Installation & Setup
    slide, content_slide = add_content_slide(prs, "Installation & Setup")
    tf = content_slide.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Step 1: Install Python 3.x"
    p.font.size = Pt(13)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Download and install Python from python.org"
    p.level = 1
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "Step 2: Clone/Download Project"
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Get the project files from the repository"
    p.level = 1
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "Step 3: Install Dependencies"
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "pip install -r requirements.txt"
    p.level = 1
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "Step 4: Run Application"
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "python timetable_gui.py or run.bat"
    p.level = 1
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    
    # SLIDE 20: Use Cases
    slide, content_box = add_content_slide(prs, "Use Cases")
    tf = content_box.text_frame
    tf.clear()
    
    use_cases = [
        "School Schedule Management",
        "University Time Table Preparation",
        "College Classroom Allocation",
        "Department-wise Planning",
        "Semester Course Scheduling",
        "Lab & Practical Slot Assignment",
        "Teacher Load Distribution",
        "Conflict Resolution & Detection",
        "Academic Calendar Planning",
        "Resource Optimization",
    ]
    
    for i, use_case in enumerate(use_cases):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"→ {use_case}"
        p.level = 0
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # SLIDE 21: Future Enhancements
    slide, content_box = add_content_slide(prs, "Future Enhancements")
    tf = content_box.text_frame
    tf.clear()
    
    enhancements = [
        "Web-based Interface (Django/Flask)",
        "Mobile Application Support",
        "Automated Conflict Detection",
        "Email Notifications",
        "Real-time Synchronization",
        "Advanced Analytics & Reporting",
        "Multi-institution Support",
        "AI-based Schedule Optimization",
        "Integration with Student System",
        "Cloud Deployment Options",
    ]
    
    for i, enhancement in enumerate(enhancements):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"▶ {enhancement}"
        p.level = 0
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # SLIDE 22: Conclusion
    slide, content_box = add_content_slide(prs, "Conclusion")
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "The Time Table Management System provides a complete, scalable solution for educational institutions."
    p.font.size = Pt(14)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "With its robust database design, intuitive GUI, and comprehensive query capabilities, it streamlines schedule creation and management."
    p.font.size = Pt(14)
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "The system is ready for immediate deployment and can be easily extended with additional features."
    p.font.size = Pt(14)
    p.space_before = Pt(15)
    
    # SLIDE 23: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)
    
    qa_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    qa_frame = qa_box.text_frame
    qa_frame.word_wrap = True
    p = qa_frame.paragraphs[0]
    p.text = "QUESTIONS?"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 133, 244)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('TimeTable_Management_System_Presentation.pptx')
    print(f"✅ PowerPoint presentation generated: TimeTable_Management_System_Presentation.pptx")

if __name__ == "__main__":
    generate_presentation()
